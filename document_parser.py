import os
import zipfile
import xml.etree.ElementTree as ET

# Importaciones condicionales para evitar fallos si las dependencias aún se están instalando
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

def extract_text_from_txt_or_md(file_path):
    encodings = ['utf-8', 'latin-1', 'cp1252']
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    raise ValueError(f"No se pudo decodificar el archivo de texto: {file_path}")

def extract_text_from_docx_fallback(file_path):
    """Método de fallback usando zipfile estándar en caso de que python-docx falle."""
    try:
        doc = zipfile.ZipFile(file_path)
        xml_content = doc.read('word/document.xml')
        root = ET.fromstring(xml_content)
        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
        paragraphs = []
        for paragraph in root.iter('{' + ns['w'] + '}p'):
            texts = [node.text for node in paragraph.iter('{' + ns['w'] + '}t') if node.text]
            paragraphs.append("".join(texts) if texts else "")
        return "\n".join(paragraphs)
    except Exception as e:
        raise ValueError(f"Fallo en el fallback de extracción docx: {e}")

def extract_text_from_docx(file_path):
    if Document is None:
        # Si python-docx no está instalado aún, usamos el fallback
        return extract_text_from_docx_fallback(file_path)
    
    try:
        doc = Document(file_path)
        full_text = []
        
        # 1. Extraer texto de párrafos normales
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        # 2. Extraer texto de TABLAS (Crítico para minutas con tablas de acuerdos)
        for table in doc.tables:
            for row in table.rows:
                # Filtrar celdas consecutivas idénticas para evitar duplicados por celdas fusionadas
                row_text = []
                last_cell_text = None
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text and cell_text != last_cell_text:
                        row_text.append(cell_text)
                        last_cell_text = cell_text
                if row_text:
                    full_text.append(" | ".join(row_text))
                    
        return "\n".join(full_text)
    except Exception as e:
        return extract_text_from_docx_fallback(file_path)

def extract_text_from_pdf(file_path):
    if PdfReader is None:
        raise ImportError("La librería 'pypdf' no está instalada. Por favor ejecuta 'pip install pypdf'")
    
    try:
        reader = PdfReader(file_path)
        text = []
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"Error al leer PDF {file_path}: {e}")

def extract_text_from_pptx(file_path):
    if Presentation is None:
        raise ImportError("La librería 'python-pptx' no está instalada. Por favor ejecuta 'pip install python-pptx'")
    
    try:
        prs = Presentation(file_path)
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Diapositiva {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if len(slide_text) > 1:
                text.append("\n".join(slide_text))
        return "\n\n".join(text)
    except Exception as e:
        raise ValueError(f"Error al leer PPTX {file_path}: {e}")

def parse_document(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"El archivo no existe: {file_path}")
        
    _, ext = os.path.splitext(file_path.lower())
    
    if ext in ['.txt', '.md']:
        return extract_text_from_txt_or_md(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    else:
        raise NotImplementedError(f"Formato de archivo no soportado: {ext}")

if __name__ == "__main__":
    # Script de prueba local
    import sys
    if len(sys.argv) > 1:
        ruta = sys.argv[1]
        try:
            print(f"\n=========================================")
            print(f"PROBANDO EXTRACCIÓN DE: {os.path.basename(ruta)}")
            print(f"=========================================\n")
            contenido = parse_document(ruta)
            # Imprimir los primeros 1500 caracteres
            print(contenido[:1500])
            if len(contenido) > 1500:
                print("\n[... Contenido truncado por visualización ...]")
            print(f"\n=========================================")
            print(f"¡ÉXITO! Longitud total: {len(contenido)} caracteres.")
            print(f"=========================================\n")
        except Exception as e:
            print(f"Error durante la prueba: {e}")
    else:
        print("Uso: python document_parser.py <ruta_del_archivo>")
